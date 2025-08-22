#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML到PPTX转换器
将templates.html转换为PowerPoint演示文稿
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import requests
from urllib.parse import urlparse

class HTMLToPPTXConverter:
    def __init__(self):
        self.prs = Presentation()
        self.slide_width = Inches(13.33)  # 标准16:9宽屏
        self.slide_height = Inches(7.5)
        
        # 颜色定义 (从HTML CSS中提取)
        self.colors = {
            'purple': RGBColor(149, 66, 157),     # #95429D
            'dark_purple': RGBColor(99, 32, 113), # #632071
            'blue': RGBColor(89, 203, 232),       # #59CBE8
            'green': RGBColor(164, 200, 89),      # #a4c859
            'dark_blue': RGBColor(0, 40, 85),     # #002855
            'white': RGBColor(255, 255, 255)
        }
    
    def download_logo(self, url, filename):
        """下载Logo图片"""
        try:
            response = requests.get(url)
            if response.status_code == 200:
                with open(filename, 'wb') as f:
                    f.write(response.content)
                return filename
        except Exception as e:
            print(f"无法下载Logo: {e}")
        return None
    
    def create_title_slide(self):
        """创建标题幻灯片"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景图片
        bg_path = "/Users/j.h.huang/Desktop/ppt_Solstice/templates/bg.png"
        if os.path.exists(bg_path):
            slide.shapes.add_picture(bg_path, 0, 0, self.slide_width, self.slide_height)
        
        # 添加标题头部
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(4), Inches(1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.colors['purple']
        header_shape.line.fill.background()
        
        # 设置标题文本
        header_text = header_shape.text_frame
        header_text.text = "FRESCO Study"
        header_text.paragraphs[0].font.size = Pt(20)
        header_text.paragraphs[0].font.color.rgb = self.colors['blue']
        header_text.paragraphs[0].font.bold = True
        header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def create_main_content_slide(self):
        """创建主要内容幻灯片"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景图片
        bg_path = "/Users/j.h.huang/Desktop/ppt_Solstice/templates/bg.png"
        if os.path.exists(bg_path):
            slide.shapes.add_picture(bg_path, 0, 0, self.slide_width, self.slide_height)
        
        # 添加标题头部
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(4), Inches(1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = self.colors['purple']
        header_shape.line.fill.background()
        
        header_text = header_shape.text_frame
        header_text.text = "FRESCO Study"
        header_text.paragraphs[0].font.size = Pt(20)
        header_text.paragraphs[0].font.color.rgb = self.colors['blue']
        header_text.paragraphs[0].font.bold = True
        header_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 下载并添加Logo
        logo_url = "https://www.fruzaqlahcp.com/sites/default/files/2024-04/fruzaqla_logo_r_rgb.png"
        logo_path = "/Users/j.h.huang/Desktop/ppt_Solstice/temp/fruzaqla_logo.png"
        
        # 确保temp目录存在
        os.makedirs(os.path.dirname(logo_path), exist_ok=True)
        
        if self.download_logo(logo_url, logo_path):
            logo_shape = slide.shapes.add_picture(
                logo_path, 
                Inches(9), Inches(1.5), 
                width=Inches(3)
            )
        
        # 添加侧边栏内容
        sidebar_left = Inches(8.5)
        sidebar_top = Inches(3)
        sidebar_width = Inches(4.5)
        
        # 重要安全信息标题
        title_shape = slide.shapes.add_textbox(
            sidebar_left, sidebar_top, 
            sidebar_width, Inches(0.5)
        )
        title_frame = title_shape.text_frame
        title_frame.text = "IMPORTANT SAFETY INFORMATION (continued)"
        title_frame.paragraphs[0].font.size = Pt(14)
        title_frame.paragraphs[0].font.color.rgb = self.colors['purple']
        title_frame.paragraphs[0].font.bold = True
        
        # 警告和注意事项副标题
        subtitle_shape = slide.shapes.add_textbox(
            sidebar_left, Inches(3.7), 
            sidebar_width, Inches(0.4)
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = "WARNINGS AND PRECAUTIONS (continued)"
        subtitle_frame.paragraphs[0].font.size = Pt(12)
        subtitle_frame.paragraphs[0].font.color.rgb = self.colors['dark_blue']
        subtitle_frame.paragraphs[0].font.bold = True
        
        # 主要内容文本
        content_shape = slide.shapes.add_textbox(
            sidebar_left, Inches(4.2), 
            sidebar_width, Inches(2)
        )
        content_frame = content_shape.text_frame
        content_frame.word_wrap = True
        content_text = ("Impaired Wound Healing. In 911 patients with mCRC treated with FRUZAQLA, "
                       "1 patient experienced a Grade 2 event of wound dehiscence. Do not administer "
                       "FRUZAQLA for at least 2 weeks prior to major surgery. Do not administer FRUZAQLA "
                       "for at least 2 weeks after major surgery and until adequate wound healing. "
                       "The safety of resumption of FRUZAQLA after resolution of wound healing "
                       "complications has not been established.")
        content_frame.text = content_text
        content_frame.paragraphs[0].font.size = Pt(10)
        content_frame.paragraphs[0].font.color.rgb = self.colors['dark_blue']
        
        # 底部重要信息
        footer_shape = slide.shapes.add_textbox(
            sidebar_left, Inches(6.5), 
            sidebar_width, Inches(0.8)
        )
        footer_frame = footer_shape.text_frame
        footer_frame.word_wrap = True
        footer_text = ("Please see additional Important Safety Information throughout, "
                      "full Important Safety Information, and full Prescribing Information for FRUZAQLA.")
        footer_frame.text = footer_text
        footer_frame.paragraphs[0].font.size = Pt(10)
        footer_frame.paragraphs[0].font.color.rgb = self.colors['dark_blue']
        footer_frame.paragraphs[0].font.bold = True
        
        return slide
    
    def create_navigation_slide(self):
        """创建导航幻灯片（模拟HTML底部的导航栏）"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景
        bg_path = "/Users/j.h.huang/Desktop/ppt_Solstice/templates/bg.png"
        if os.path.exists(bg_path):
            slide.shapes.add_picture(bg_path, 0, 0, self.slide_width, self.slide_height)
        
        # 导航项目
        nav_items = [
            ("HOME", "home"),
            ("UNMET\nNEED/MOA", "purple"),
            ("FRESCO-2", "green"),
            ("EFFICACY:\nOS", "purple"),
            ("EFFICACY:\nPFS", "purple"),
            ("SAFETY", "purple"),
            ("FRESCO", "blue"),
            ("EFFICACY:\nOS", "purple"),
            ("EFFICACY:\nPFS", "purple"),
            ("SAFETY", "purple"),
            ("QOL/\nDOSING", "purple"),
            ("ASSIST/\nISI", "purple"),
            ("ISI/\nREFERENCES", "purple")
        ]
        
        # 计算导航按钮的位置
        box_width = Inches(1)
        box_height = Inches(1)
        start_x = Inches(0.5)
        start_y = Inches(6)
        
        for i, (text, color_type) in enumerate(nav_items):
            x_pos = start_x + (i * Inches(1.02))
            
            # 创建按钮形状
            if text == "HOME":
                # 添加Home图标
                home_icon_path = "/Users/j.h.huang/Desktop/ppt_Solstice/templates/HomeIcon.png"
                if os.path.exists(home_icon_path):
                    button_shape = slide.shapes.add_picture(
                        home_icon_path, x_pos, start_y, 
                        width=Inches(0.6), height=Inches(0.6)
                    )
            else:
                # 创建文本按钮
                button_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 
                    x_pos, start_y, box_width, box_height
                )
                
                # 设置按钮颜色
                if color_type == "purple":
                    button_shape.fill.solid()
                    button_shape.fill.fore_color.rgb = self.colors['purple']
                elif color_type == "green":
                    button_shape.fill.solid()
                    button_shape.fill.fore_color.rgb = self.colors['green']
                elif color_type == "blue":
                    button_shape.fill.solid()
                    button_shape.fill.fore_color.rgb = RGBColor(61, 161, 187)
                
                button_shape.line.fill.background()
                
                # 添加文本
                text_frame = button_shape.text_frame
                text_frame.text = text
                text_frame.paragraphs[0].font.size = Pt(8)
                text_frame.paragraphs[0].font.color.rgb = self.colors['white']
                text_frame.paragraphs[0].font.bold = True
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = 1  # 垂直居中
        
        return slide
    
    def convert(self):
        """执行转换"""
        print("开始转换HTML到PPTX...")
        
        # 设置幻灯片尺寸为16:9
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        
        # 创建幻灯片
        self.create_title_slide()
        self.create_main_content_slide()
        self.create_navigation_slide()
        
        # 保存PPTX文件
        output_path = "/Users/j.h.huang/Desktop/ppt_Solstice/output/fresco_study_presentation.pptx"
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        self.prs.save(output_path)
        print(f"转换完成！PPTX文件已保存至: {output_path}")
        
        return output_path

if __name__ == "__main__":
    converter = HTMLToPPTXConverter()
    converter.convert() 